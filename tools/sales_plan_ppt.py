#!/usr/bin/env python3
"""sales_plan_ppt.py — Populate Adobe Lite DX Business Plan Template from content_map.json.

Reads a ``content_map.json`` produced by the ``/sales plan`` Claude Code skill
and populates a cloned copy of the Adobe Lite DX Business Plan PowerPoint
template. All intelligence work happens upstream in Claude; this script is
purely mechanical: read JSON, locate shapes by name, rebuild text frames.

Usage::

    py -3 sales_plan_ppt.py \\
        --content-map  "C:/Users/arjaiswa/Desktop/claude-workspace/deals/company/content/content_map.json" \\
        --template     "C:/Users/arjaiswa/OneDrive - Adobe/Lite DX Business Plan Template - Arpit.pptx" \\
        --output       "C:/Users/arjaiswa/Desktop/claude-workspace/deals/company/artifacts"

Dependencies:
    python-pptx >= 0.6.23
    lxml >= 4.0

Slides populated (shape-name lookup unless noted):
    1   Company name title ("Company — DALP Account Plan") + company logo
    2   Executive summary table — business issue, big idea, objectives,
        challenges, Adobe differentiated solution
    3   Company overview — background (structured 10-line format), LOB,
        account intel, Adobe strengths, opportunities
    4   Performance history — upcoming renewals, renewal strategy
        (from CLI args or Clari/Panorama; [FILL] placeholder otherwise)
    5   Market landscape — trends, goals, digital priorities, opportunities,
        challenges, partner strategy, implementation partners
    6   Buying committee table — ALL non-placeholder shapes cleared first;
        fresh 5-col table inserted (Name, Title, Role, Reason, Attitude)
    7   Value strategy — big idea, tagline, business issue, portfolio plays,
        path to value, like customers
    8   FY opportunities — pipeline summary
        (from CLI args or Clari; [FILL] placeholder otherwise)
    9   FY timeline — monthly touchpoints distributed across H1 (Dec–Jul)
        and H2 (Aug–Nov) columns; backward-compatible with flat text format
    11  Big Idea appendix — goals, challenges, initiatives, impact rows,
        full Big Idea paragraph, tagline

Notes:
    - The source template is never modified; the script always clones it.
    - If an output file already exists it is renamed with a ``_HHMM`` suffix
      before the new file is written, preserving prior versions.
    - Shape names are matched exactly (case-sensitive). Run the shape
      inspector snippet in README.md if slides appear blank after population.

Version: 1.3.0
Changelog: https://github.com/arpitjaiswal-0701/sales-plan-skill/blob/master/CHANGELOG.md
"""

import argparse
import copy
import json
import re
import shutil
import sys
import time
import urllib.request
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Set

# Windows console may be cp1252 — reconfigure to utf-8 so Unicode output works
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree


__version__ = "1.3.0"

TEMPLATE_DEFAULT = r"C:\Users\arjaiswa\OneDrive - Adobe\Lite DX Business Plan Template - Arpit.pptx"


# ── Shape helpers ─────────────────────────────────────────────────────────────

def find_shape(container: Any, name: str) -> Optional[Any]:
    """Recursively search a slide or group for a shape with the given name.

    Args:
        container: A python-pptx Slide or GroupShapes object with a ``.shapes``
            collection. Any other type is treated as having no shapes.
        name: The exact shape name to match (case-sensitive).

    Returns:
        The first matching shape, or ``None`` if no match is found.
    """
    shapes = getattr(container, 'shapes', [])
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            result = find_shape(shape, name)
            if result:
                return result
    return None


def find_shapes_batch(container: Any, names: Set[str]) -> Dict[str, Any]:
    """Single-pass recursive search for multiple shape names.

    More efficient than calling find_shape() N times on the same slide because
    the shape tree (including nested GROUP containers) is traversed only once
    regardless of how many names are requested. On slides with several group
    shapes (e.g. slide 7 with Group 218 → Group 60) this avoids redundant
    recursion.

    Args:
        container: A python-pptx Slide or GroupShapes object.
        names: Set of exact shape names to collect.

    Returns:
        Dict mapping found shape names to shape objects. Names not found in the
        slide are omitted from the result.
    """
    result: Dict[str, Any] = {}
    _collect_shapes(container, names, result)
    return result


def _collect_shapes(container: Any, names: Set[str], result: Dict[str, Any]) -> None:
    """Recursive helper for find_shapes_batch."""
    if len(result) >= len(names):
        return  # early exit — all targets found
    for shape in getattr(container, 'shapes', []):
        if shape.name in names and shape.name not in result:
            result[shape.name] = shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            _collect_shapes(shape, names, result)


def _normalize_txbody(txBody: Any) -> None:
    """Apply anchor=top and normAutofit to an ``<a:txBody>`` element directly.

    Shared implementation called by both :func:`normalize_body_props`
    (shape-level entry point) and :func:`set_cell_text` (table cell entry
    point), guaranteeing identical body-property behavior across all text
    frames without code duplication.

    Args:
        txBody: An lxml element representing an ``<a:txBody>``.
    """
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    bodyPr.set('anchor', 't')
    bodyPr.attrib.pop('anchorCtr', None)
    for tag in (qn('a:noAutofit'), qn('a:spAutoFit')):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)
    if not bodyPr.findall(qn('a:normAutofit')):
        etree.SubElement(bodyPr, qn('a:normAutofit'))


def set_text(shape: Any, text: str) -> None:
    """Populate a shape's text frame, replacing all existing content.

    Extracts the first run's ``<a:rPr>`` (font size, bold, color, typeface)
    before clearing the text frame, then rebuilds one ``<a:p>`` per
    newline-delimited line. This guarantees no residual template text or
    field elements survive in the XML.

    Newline splitting enables multi-line bullet content from
    ``content_map.json`` (e.g. ``"• Trend 1\\n• Trend 2"``) to render as
    separate paragraphs rather than a single run-together line.

    Calls :func:`normalize_body_props` after writing to fix vertical anchor
    and autofit settings on the shape.

    Args:
        shape: A python-pptx shape with a text frame. Silently no-ops if
            the shape is falsy or does not have a text frame.
        text: Content to write. ``\\n`` characters produce separate
            ``<a:p>`` elements; an empty string produces one blank paragraph.
    """
    if not shape or not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn('a:p'))
    if not paras:
        return

    # Extract rPr from first run before clearing — captures font, size, color,
    # theme references, and any custom run properties set by the template.
    rPr: Optional[Any] = None
    runs = paras[0].findall(qn('a:r'))
    if runs:
        rPr_elem = runs[0].find(qn('a:rPr'))
        if rPr_elem is not None:
            rPr = copy.deepcopy(rPr_elem)

    for para in paras:
        txBody.remove(para)

    for line in (text.split('\n') if text else ['']):
        para = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(para, qn('a:pPr'))
        pPr.set('algn', 'l')
        r = etree.SubElement(para, qn('a:r'))
        if rPr is not None:
            r.insert(0, copy.deepcopy(rPr))
        t = etree.SubElement(r, qn('a:t'))
        t.text = line

    normalize_body_props(shape)


def normalize_body_props(shape: Any) -> None:
    """Fix text frame body properties that cause alignment and clipping defects.

    Sets three properties on ``<a:bodyPr>`` that are frequently misconfigured
    in PowerPoint templates and produce three distinct visual failures:

    - ``anchor="t"`` — aligns text to the top of the shape.
    - ``<a:normAutofit>`` — shrinks the font proportionally to fit the shape's
      fixed dimensions, replacing ``<a:noAutofit>`` (clips overflow) and
      ``<a:spAutoFit>`` (resizes the shape).
    - Removes ``anchorCtr`` — eliminates secondary vertical centering.

    Delegates to :func:`_normalize_txbody` for the actual XML manipulation.

    Args:
        shape: A python-pptx shape with a text frame. Silently no-ops if the
            shape is falsy, lacks a text frame, or has no ``<a:bodyPr>``.
    """
    if not shape or not shape.has_text_frame:
        return
    _normalize_txbody(shape.text_frame._txBody)


def set_cell_text(cell: Any, text: str) -> None:
    """Write text into a table cell using XML-level rebuild with normAutofit.

    Fully rebuilds all ``<a:p>`` elements inside the cell's ``<a:txBody>``,
    preserving the first run's ``<a:rPr>`` formatting. Delegates body-property
    normalization to :func:`_normalize_txbody` to prevent row-height expansion
    and text clipping. Supports ``\\n``-delimited multi-line text.

    Args:
        cell: A python-pptx table cell object.
        text: The text string to write into the cell.
    """
    try:
        tf = cell.text_frame
        txBody = tf._txBody

        # Preserve run formatting from first existing paragraph
        rPr: Optional[Any] = None
        existing_paras = txBody.findall(qn('a:p'))
        if existing_paras:
            runs = existing_paras[0].findall(qn('a:r'))
            if runs:
                rPr_elem = runs[0].find(qn('a:rPr'))
                if rPr_elem is not None:
                    rPr = copy.deepcopy(rPr_elem)

        # Remove all existing paragraphs
        for para in existing_paras:
            txBody.remove(para)

        # Rebuild one paragraph per line
        for line in (text.split('\n') if text else ['']):
            para = etree.SubElement(txBody, qn('a:p'))
            pPr = etree.SubElement(para, qn('a:pPr'))
            pPr.set('algn', 'l')
            r = etree.SubElement(para, qn('a:r'))
            if rPr is not None:
                r.insert(0, copy.deepcopy(rPr))
            t = etree.SubElement(r, qn('a:t'))
            t.text = line

        _normalize_txbody(txBody)
    except Exception:
        pass


def get_tables(slide: Any) -> List[Any]:
    """Return all table shapes on a slide in document order.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        List of shapes where ``shape_type == 19`` (``MSO_SHAPE_TYPE.TABLE``).
        Returns an empty list if the slide has no tables.
    """
    return [s for s in slide.shapes if s.shape_type == 19]


def fetch_logo(domain: str, output_dir: Path) -> Optional[str]:
    """Download company logo from Clearbit Logo API.

    Args:
        domain: Company domain, e.g. ``"appliedmaterials.com"``.
        output_dir: Directory to save the cached logo file.

    Returns:
        Absolute path string of the saved logo PNG, or ``None`` if download
        fails (network unavailable, domain not found, or any exception).
    """
    if not domain:
        return None
    logo_path = output_dir / "logo_cache.png"
    if logo_path.exists():
        return str(logo_path)
    url = f"https://logo.clearbit.com/{domain}"
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = resp.read()
        with open(logo_path, 'wb') as f:
            f.write(data)
        return str(logo_path)
    except Exception:
        return None


def _apply_cell_style(cell: Any, bold: bool = False, white_text: bool = False,
                      bg_hex: Optional[str] = None) -> None:
    """Apply font size and optional bold/color/background to a table cell.

    Called after :func:`set_cell_text` since that function rebuilds all runs
    from scratch. Sets font size to 8.5pt on all runs; optionally applies
    bold, white text colour, and a solid background fill on the cell.

    Args:
        cell: A python-pptx table cell.
        bold: Set run bold attribute when ``True``.
        white_text: Apply white (#FFFFFF) font colour when ``True``.
        bg_hex: Six-character hex string for cell background (e.g. ``'FA0F00'``).
    """
    txBody = cell.text_frame._txBody
    for r_el in txBody.iter(qn('a:r')):
        rPr = r_el.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.SubElement(r_el, qn('a:rPr'))
        rPr.set('sz', '850')  # 8.5pt in hundredths of a point
        if bold:
            rPr.set('b', '1')
        if white_text:
            for existing in rPr.findall(qn('a:solidFill')):
                rPr.remove(existing)
            solidFill = etree.SubElement(rPr, qn('a:solidFill'))
            etree.SubElement(solidFill, qn('a:srgbClr')).set('val', 'FFFFFF')
    if bg_hex:
        tcPr = cell._tc.get_or_add_tcPr()
        for existing in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(existing)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr')).set('val', bg_hex)


# ── Slide populators ──────────────────────────────────────────────────────────

def slide_1(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 1 — DALP title and company logo.

    Sets the title placeholder to "{Company Name} — DALP Account Plan" and
    inserts the company logo into the lower-right corner of the slide.
    Checks ctx['logo_path'] first (pre-fetched by main); falls back to a
    direct Clearbit download if not present in ctx. Silently skips logo if
    download fails.
    """
    company = cm.get('company_name', '[FILL: Company Name]')
    title_shape = find_shape(slide, 'Title 1')
    if title_shape:
        set_text(title_shape, f"{company} — DALP Account Plan")

    # Use pre-fetched logo from ctx when available; fall back to direct fetch
    logo_path = ctx.get('logo_path') if ctx else None
    if not logo_path:
        domain = cm.get('company_domain', '')
        if domain and ctx:
            output_dir = Path(ctx.get('output_dir', '.'))
            logo_path = fetch_logo(domain, output_dir)

    if logo_path:
        try:
            # Lower-right corner: 2.5" wide × 1.5" tall, 0.4" inset from edges
            slide.shapes.add_picture(
                logo_path,
                Inches(10.43), Inches(5.60), Inches(2.5), Inches(1.5)
            )
            print("  ↳ Logo inserted")
        except Exception as e:
            print(f"  ↳ Logo insert failed: {e}")


def slide_2(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 2 — Executive Summary table and last-refreshed date.

    The executive summary table (Table 5) has 8 columns; content lands in
    row 1 at the anchor cells for each merged column group:
      col 0 — Business Issue (single)
      col 1 — Big Idea (spans cols 1-2)
      col 3 — Company Objectives (spans cols 3-4)
      col 5 — Challenges (spans cols 5-6)
      col 7 — Adobe Differentiated Solution (single)
    """
    s = cm.get('slide_2', {})
    tables = get_tables(slide)
    if tables:
        t = tables[0].table
        mapping: Dict[int, str] = {
            0: s.get('business_issue',    '[FILL: Business Issue]'),
            1: s.get('big_idea',           '[FILL: Big Idea]'),
            3: s.get('company_objectives', '[FILL: Company Objectives]'),
            5: s.get('challenges',         '[FILL: Key Challenges]'),
            7: s.get('adobe_solution',     '[FILL: Adobe Differentiated Solution]'),
        }
        for col, text in mapping.items():
            try:
                set_cell_text(t.cell(1, col), text)
            except Exception:
                pass

    date_shape = find_shape(slide, 'TextBox 15')
    if date_shape:
        set_text(date_shape, f"Last Refreshed on: {cm.get('date', datetime.today().strftime('%Y-%m-%d'))}")


def slide_3(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 3 — Company Overview content boxes.

    Rectangle 11 (LOB A / main account overview) and Rectangle 32 (LOB B /
    department overview) both expect a 10-line structured format matching the
    template's labeled fields. Content should be newline-delimited with each
    line starting "FieldName: value".

    Uses find_shapes_batch for a single-pass shape lookup across all 5 targets.
    """
    s = cm.get('slide_3', {})
    targets: Dict[str, str] = {
        'Rectangle 11': s.get('account_background',    '[FILL: Annual Revenue:\nOnline Revenue:\nSize:\nBusiness Focus & Major Divisions:\nIndustry & Position:\nGeneral digital maturity:\nGeography:\nCompetitive Footprint:\nAzure or AWS Commitments:\nPartner Footprint:]'),
        'Rectangle 32': s.get('account_background_lob','[FILL: Annual Revenue:\nOnline Revenue:\nSize:\nBusiness Focus & Major Divisions:\nIndustry & Position:\nGeneral digital maturity:\nGeography:\nCompetitive Footprint:\nAzure or AWS Commitments:\nPartner Footprint:]'),
        'Rectangle 34': s.get('account_intel',          '[FILL: Account intelligence — whitespace rationale, Tier 1 reasoning]'),
        'Rectangle 51': s.get('adobe_strengths',        '[FILL: Adobe advantages and capabilities for this account]'),
        'Rectangle 59': s.get('opportunities',          '[FILL: Specific opportunity areas to create value]'),
    }
    found = find_shapes_batch(slide, set(targets.keys()))
    for name, text in targets.items():
        shape = found.get(name)
        if shape:
            set_text(shape, text)


def slide_4(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 4 — Performance History (renewals and renewal strategy).

    Populated from CLI arguments (``--arr``, ``--renewal``) if provided by the
    ``/sales plan`` skill; otherwise a ``[FILL]`` placeholder is written.
    Complete population requires Clari/Panorama data.

    Template table layout: index 0 = instructions, 1 = Upcoming Renewals,
    2 = Renewal Strategy.
    """
    s = cm.get('slide_4', {})
    tables = get_tables(slide)
    if len(tables) > 1:
        try:
            set_cell_text(tables[1].table.cell(1, 0),
                          s.get('upcoming_renewals',
                                '[FILL: Upcoming renewals — Solution | $Value | Renewal Date | Driver. Source: Clari/Panorama.]'))
        except Exception:
            pass
    if len(tables) > 2:
        try:
            set_cell_text(tables[2].table.cell(1, 0),
                          s.get('renewal_strategy',
                                '[FILL: Renewal strategy and risk. Source: Panorama / Clari.]'))
        except Exception:
            pass


def slide_5(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 5 — Market Landscape (trends, goals, priorities, partners).

    Uses find_shapes_batch for a single-pass shape lookup across all 7 targets.
    """
    s = cm.get('slide_5', {})
    targets: Dict[str, str] = {
        'Rectangle 49':  s.get('market_trends',          '[FILL: Broad industry/market developments — 3-4 bullets]'),
        'Rectangle 9':   s.get('company_goals',           '[FILL: Company strategic goals — 3-4 bullets]'),
        'Rectangle 15':  s.get('digital_priorities',      '[FILL: Digital priorities and technology strategy]'),
        'Rectangle 14':  s.get('market_opportunities',    '[FILL: Specific areas to create value or solve problems]'),
        'Rectangle 4':   s.get('customer_challenges',     '[FILL: Key obstacles customer must overcome — in customer language]'),
        'Rectangle 183': s.get('partner_strategy',        '[FILL: Partner strategy and relationships]'),
        'Rectangle 187': s.get('implementation_partners', '[FILL: Historical implementation partners]'),
    }
    found = find_shapes_batch(slide, set(targets.keys()))
    for name, text in targets.items():
        shape = found.get(name)
        if shape:
            set_text(shape, text)


def slide_6(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 6 — Buying Committee table.

    Removes ALL non-placeholder shapes from the slide (org chart boxes,
    connector lines, legend text boxes, partner relationship rectangles,
    photo ovals, and all GROUP shapes) before inserting the new table. This
    cleanly replaces the template's pre-built org chart diagram.

    Only the title placeholder ('Customer Organization Chart') is preserved.

    Table layout (5 columns, 12.5" total):
        Name (2.0") | Title (2.5") | Role (1.8") | Reason for Engagement (4.2") | Attitude (2.0")
        Header row: Adobe red background (#FA0F00), white bold text at 8.5pt
        Data rows: up to 10 contacts from content_map.json slide_6.buying_committee

    Uses set_cell_text + _apply_cell_style (XML-level) for all cells, ensuring
    normAutofit and consistent body properties on every row.
    """
    # Remove ALL non-placeholder shapes — org chart, legend, lines, pictures, groups
    to_remove = [s for s in slide.shapes if not s.is_placeholder]
    for s in to_remove:
        s._element.getparent().remove(s._element)

    contacts = cm.get('slide_6', {}).get('buying_committee', [])
    if not contacts:
        return

    rows = min(len(contacts), 10) + 1  # +1 for header row
    tbl_shape = slide.shapes.add_table(
        rows, 5, Inches(0.3), Inches(1.15), Inches(12.5), Pt(20) * rows
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(2.5)
    tbl.columns[2].width = Inches(1.8)
    tbl.columns[3].width = Inches(4.2)
    tbl.columns[4].width = Inches(2.0)

    headers = ['Name', 'Title', 'Role', 'Reason for Engagement', 'Attitude toward Adobe']
    for ci, header_text in enumerate(headers):
        set_cell_text(tbl.cell(0, ci), header_text)
        _apply_cell_style(tbl.cell(0, ci), bold=True, white_text=True, bg_hex='FA0F00')

    fields = ['name', 'title', 'role', 'reason', 'attitude']
    for ri, contact in enumerate(contacts[:10]):
        for ci, key in enumerate(fields):
            set_cell_text(tbl.cell(ri + 1, ci), str(contact.get(key, '')))
            _apply_cell_style(tbl.cell(ri + 1, ci))


def slide_7(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 7 — Value Strategy.

    Populates big idea, tagline, business issue, portfolio plays, path to
    value, and like-customer references. The ``Rectangle 63`` (like customers)
    shape is nested two levels deep in Group 218 → Group 60 and is found via
    the recursive find_shapes_batch single-pass traversal.
    """
    s = cm.get('slide_7', {})
    targets: Dict[str, str] = {
        'Rectangle 25':  s.get('big_idea',       '[FILL: Big Idea in customer language — 2-3 sentences]'),
        'Rectangle 26':  s.get('tagline',          '[FILL: One-sentence tagline]'),
        'Rectangle 69':  s.get('business_issue',   '[FILL: Business issue impacting performance and goals]'),
        'Rectangle 84':  s.get('portfolio_plays',  '[FILL: Adobe Portfolio / Sales Plays to pitch]'),
        'Rectangle 217': s.get('path_to_value',    '[FILL: Path to Value and Budget alignment. Potential pipeline: $X]'),
        'Rectangle 63':  s.get('like_customers',   '[FILL: Peer companies with similar ALM/Adobe DX deployments]'),
    }
    found = find_shapes_batch(slide, set(targets.keys()))
    for name, text in targets.items():
        shape = found.get(name)
        if shape:
            set_text(shape, text)


def slide_8(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 8 — FY Opportunities summary.

    Populated from CLI arguments (``--arr``, ``--stage``, ``--close-date``) if
    provided; otherwise a ``[FILL]`` placeholder directs the user to Clari.
    """
    s = cm.get('slide_8', {})
    tables = get_tables(slide)
    if tables:
        try:
            set_cell_text(tables[0].table.cell(0, 0),
                          s.get('opportunities_summary',
                                '[FILL: Complete in Clari. Paste pipeline table: Opp Name | Solution | ARR | Stage | Close Date | Next Step.]'))
        except Exception:
            pass


def slide_9(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 9 — FY Timeline monthly touchpoints.

    Supports two content_map.json formats:

    **New format** (preferred): nested objects with monthly keys distribute
    touchpoints across individual month columns:
    ``slide_9.h1 = {"december": "...", "january": "...", ...}``
    ``slide_9.h2 = {"august": "...", ..., "november": "..."}``

    **Legacy format** (backward-compatible): flat strings written to the
    first month column of each half-year table:
    ``slide_9.touchpoints_h1`` and ``slide_9.touchpoints_h2``

    H1 table covers December (col 1) through July (col 8).
    H2 table covers August (col 1) through November (col 4).
    """
    s = cm.get('slide_9', {})
    tables = get_tables(slide)

    H1_MONTHS = ['december', 'january', 'february', 'march',
                 'april',    'may',     'june',     'july']
    H2_MONTHS = ['august', 'september', 'october', 'november']

    if tables:
        t = tables[0].table
        h1 = s.get('h1')
        if isinstance(h1, dict):
            for col_offset, month in enumerate(H1_MONTHS):
                text = h1.get(month, '')
                if text:
                    try:
                        set_cell_text(t.cell(1, col_offset + 1), text)
                    except Exception:
                        pass
        else:
            legacy = s.get('touchpoints_h1', '')
            if legacy:
                try:
                    set_cell_text(t.cell(1, 1), legacy)
                except Exception:
                    pass

    if len(tables) > 1:
        t2 = tables[1].table
        h2 = s.get('h2')
        if isinstance(h2, dict):
            for col_offset, month in enumerate(H2_MONTHS):
                text = h2.get(month, '')
                if text:
                    try:
                        set_cell_text(t2.cell(1, col_offset + 1), text)
                    except Exception:
                        pass
        else:
            legacy = s.get('touchpoints_h2', '')
            if legacy:
                try:
                    set_cell_text(t2.cell(1, 1), legacy)
                except Exception:
                    pass


def slide_11(slide: Any, cm: Dict[str, Any], ctx: Optional[Dict] = None) -> None:
    """Populate Slide 11 — Big Idea appendix.

    Writes four labeled rows (Goals, Challenges, Initiatives, Impact) into the
    first table on the slide, then populates the full Big Idea paragraph text
    box and a standalone tagline text box.
    """
    s = cm.get('slide_11', {})
    tables = get_tables(slide)
    if tables:
        t = tables[0].table
        rows_data = [
            ('goals',       'Goals\n'),
            ('challenges',  'Challenges\n'),
            ('initiatives', 'Initiatives\n'),
            ('impact',      'Impact\n'),
        ]
        for ri, (key, prefix) in enumerate(rows_data):
            val = s.get(key, f'[FILL: {key.title()} — from company goals/priorities slide]')
            try:
                set_cell_text(t.cell(ri, 0), prefix + val)
            except Exception:
                pass

    big_idea = find_shape(slide, 'TextBox 25')
    if big_idea:
        set_text(big_idea, s.get('big_idea_paragraph', '[FILL: Big Idea paragraph in customer language]'))

    tagline = find_shape(slide, 'TextBox 11')
    if tagline:
        set_text(tagline, s.get('tagline', '[FILL: Tagline]'))


# ── Slide dispatch ─────────────────────────────────────────────────────────────

POPULATORS: Dict[int, Any] = {
    1:  slide_1,
    2:  slide_2,
    3:  slide_3,
    4:  slide_4,
    5:  slide_5,
    6:  slide_6,
    7:  slide_7,
    8:  slide_8,
    9:  slide_9,
    11: slide_11,
}


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description='Populate Lite DX Business Plan PPT from content_map.json'
    )
    parser.add_argument('--content-map', required=True, help='Path to content_map.json')
    parser.add_argument('--template', default=TEMPLATE_DEFAULT, help='Path to .pptx template')
    parser.add_argument('--output', required=True, help='Output directory path')
    args = parser.parse_args()

    with open(args.content_map, 'r', encoding='utf-8') as f:
        cm = json.load(f)

    template_path = Path(args.template)
    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    ctx: Dict[str, Any] = {'output_dir': str(output_dir)}

    # Pre-fetch logo before the populate loop so network I/O doesn't block
    # slide 1 population timing and the cached path is reused on refresh runs.
    domain = cm.get('company_domain', '')
    if domain:
        print(f"  Fetching logo for {domain}…")
        lp = fetch_logo(domain, output_dir)
        if lp:
            ctx['logo_path'] = lp
            print(f"  ↳ Cached: {Path(lp).name}")
        else:
            print(f"  ↳ Logo unavailable (network or domain not found)")

    company  = re.sub(r'[^\w\s-]', '', cm.get('company_name', 'Unknown')).strip().replace(' ', '-')
    date_str = cm.get('date', datetime.today().strftime('%Y-%m-%d'))
    out_name = f"{company}-Business-Plan-{date_str}.pptx"
    out_path = output_dir / out_name

    if out_path.exists():
        ts = datetime.now().strftime('%H%M')
        backup = output_dir / f"{company}-Business-Plan-{date_str}_{ts}.pptx"
        shutil.move(str(out_path), str(backup))
        print(f"  ↳ Prior version renamed: {backup.name}")

    shutil.copy2(str(template_path), str(out_path))
    prs = Presentation(str(out_path))

    errors: List[str] = []
    t_total = time.perf_counter()
    for i, slide in enumerate(prs.slides):
        n = i + 1
        if n in POPULATORS:
            t0 = time.perf_counter()
            try:
                POPULATORS[n](slide, cm, ctx)
                print(f"  Slide {n}: {time.perf_counter() - t0:.2f}s")
            except Exception as e:
                errors.append(f"Slide {n}: {e}")

    prs.save(str(out_path))
    print(f"\n✓ PPT saved: {out_path}  ({time.perf_counter() - t_total:.2f}s total)")

    if errors:
        print("\n⚠ Non-fatal errors:")
        for e in errors:
            print(f"  {e}")

    fill_needed: List[int] = []
    for slide_key, slide_num in [('slide_4', 4), ('slide_8', 8)]:
        vals = cm.get(slide_key, {}).values()
        if any(str(v).startswith('[FILL') for v in vals):
            fill_needed.append(slide_num)
    if fill_needed:
        print(f"\n⚑ Manual fill needed on slides: {', '.join(str(s) for s in fill_needed)}")
        print("  Source: Clari / Panorama")


if __name__ == '__main__':
    main()
